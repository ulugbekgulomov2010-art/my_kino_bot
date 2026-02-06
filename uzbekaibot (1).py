# ✅ UZBEK AI BOT (MEGA) — 🤖 AI | 🧰 Media | 📊 PPTX | 🎤 Voice | 🖼 Vision
# ------------------------------------------------------------
# pip install pyTelegramBotAPI groq ddgs requests python-pptx
# ------------------------------------------------------------

import os
import re
import json
import time
import uuid
import html
import base64
import mimetypes
import random
import sqlite3
import threading
import warnings
from datetime import datetime, timezone

import requests
import telebot
from telebot.types import (
    ReplyKeyboardMarkup, KeyboardButton,
    InlineKeyboardMarkup, InlineKeyboardButton
)

from groq import Groq
from ddgs import DDGS
from pptx import Presentation


# ================== SOZLAMALAR (API KEYLAR) ==================
# ⚠️ BU YERGA O'Z KEYLARINGIZNI QO'YING (MEN QAYTA CHIQARMAYMAN)
BOT_TOKEN = "BOT_TOKENIZNI_QO'YASIZ"
GROQ_API_KEY = "GROQ_AI DEGAN AI bor o'shani saytidan api olasiz"

CHAT_MODEL = "llama-3.1-8b-instant"
VISION_MODEL = "meta-llama/llama-4-scout-17b-16e-instruct"
WHISPER_MODEL = "whisper-large-v3"

SYSTEM_PROMPT = (
    "Sen Telegramdagi sun’iy intellekt botsan. "
    "Foydalanuvchi bilan FAQAT o‘zbek tilida gaplash. "
    "Hech qachon turkcha, ruscha yoki inglizcha yozma. "
    "Agar foydalanuvchi boshqa tilda yozsa ham, javobni o‘zbek tilida ber."
)

DB_FILE = "uzbekaibot.db"


# ================== INIT ==================
# ⚠️ MUHIM: parse_mode=None -> <urllib3...> HTML bo‘lib ketmaydi
bot = telebot.TeleBot(BOT_TOKEN, parse_mode=None)
client = Groq(api_key=GROQ_API_KEY)

session = requests.Session()
session.trust_env = False
NO_PROXY = {"http": None, "https": None}

os.makedirs("media", exist_ok=True)
os.makedirs("pptx_out", exist_ok=True)
os.makedirs("templates", exist_ok=True)


# ================== SAFE HELPERS ==================
def now_str():
    return datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M:%S")

def safe_html(s: str) -> str:
    return html.escape(str(s), quote=False)

def send_plain(chat_id: int, text: str, **kw):
    return bot.send_message(chat_id, str(text), parse_mode=None, **kw)

def reply_plain(m, text: str, **kw):
    return bot.reply_to(m, str(text), parse_mode=None, **kw)

def send_html(chat_id: int, text: str, **kw):
    return bot.send_message(chat_id, text, parse_mode="HTML", **kw)

def reply_html(m, text: str, **kw):
    return bot.reply_to(m, text, parse_mode="HTML", **kw)

def safe_error(chat_id: int, prefix: str, err: Exception):
    send_plain(chat_id, f"{prefix}: {err}")


# ================== DB (THREAD-SAFE) ==================
DB_LOCK = threading.Lock()

def db_connect():
    conn = sqlite3.connect(DB_FILE, check_same_thread=False)
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA journal_mode=WAL;")
    conn.execute("PRAGMA synchronous=NORMAL;")
    conn.execute("PRAGMA busy_timeout=7000;")
    return conn

conn = db_connect()
cur = conn.cursor()

def table_has_column(table: str, col: str) -> bool:
    rows = cur.execute(f"PRAGMA table_info({table})").fetchall()
    return any(r["name"] == col for r in rows)

def ensure_schema():
    with DB_LOCK:
        cur.execute("""
        CREATE TABLE IF NOT EXISTS users (
            user_id INTEGER PRIMARY KEY,
            chat_id INTEGER NOT NULL,
            created_at TEXT NOT NULL,
            active_ai_chat_id INTEGER
        )
        """)
        cur.execute("""
        CREATE TABLE IF NOT EXISTS ai_chats (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            name TEXT NOT NULL,
            created_at TEXT NOT NULL,
            closed_at TEXT
        )
        """)
        cur.execute("""
        CREATE TABLE IF NOT EXISTS ai_messages (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            ai_chat_id INTEGER NOT NULL,
            role TEXT NOT NULL,
            content TEXT NOT NULL,
            created_at TEXT NOT NULL
        )
        """)

        # MIGRATION
        if not table_has_column("ai_chats", "closed_at"):
            cur.execute("ALTER TABLE ai_chats ADD COLUMN closed_at TEXT")
        if not table_has_column("users", "active_ai_chat_id"):
            cur.execute("ALTER TABLE users ADD COLUMN active_ai_chat_id INTEGER")

        conn.commit()

ensure_schema()

def get_or_create_user(user_id: int, chat_id: int):
    with DB_LOCK:
        u = cur.execute("SELECT * FROM users WHERE user_id=?", (user_id,)).fetchone()
        if not u:
            cur.execute(
                "INSERT INTO users(user_id, chat_id, created_at, active_ai_chat_id) VALUES (?,?,?,NULL)",
                (user_id, chat_id, now_str())
            )
            conn.commit()
            u = cur.execute("SELECT * FROM users WHERE user_id=?", (user_id,)).fetchone()
        else:
            if u["chat_id"] != chat_id:
                cur.execute("UPDATE users SET chat_id=? WHERE user_id=?", (chat_id, user_id))
                conn.commit()
        return u

def create_ai_chat(user_id: int, name: str):
    with DB_LOCK:
        cur.execute(
            "INSERT INTO ai_chats(user_id, name, created_at, closed_at) VALUES (?,?,?,NULL)",
            (user_id, name, now_str())
        )
        cid = cur.lastrowid
        cur.execute("UPDATE users SET active_ai_chat_id=? WHERE user_id=?", (cid, user_id))
        conn.commit()
        return cid

def get_active_chat_id(user_id: int):
    with DB_LOCK:
        u = cur.execute("SELECT active_ai_chat_id FROM users WHERE user_id=?", (user_id,)).fetchone()
        return u["active_ai_chat_id"] if u else None

def set_active_chat(user_id: int, ai_chat_id: int):
    with DB_LOCK:
        cur.execute("UPDATE users SET active_ai_chat_id=? WHERE user_id=?", (ai_chat_id, user_id))
        conn.commit()

def close_chat(user_id: int, chat_id: int):
    with DB_LOCK:
        cur.execute(
            "UPDATE ai_chats SET closed_at=? WHERE id=? AND user_id=?",
            (now_str(), chat_id, user_id)
        )
        conn.commit()

def rename_chat(user_id: int, chat_id: int, new_name: str):
    with DB_LOCK:
        cur.execute("UPDATE ai_chats SET name=? WHERE id=? AND user_id=?", (new_name, chat_id, user_id))
        conn.commit()

def list_user_chats(user_id: int, limit: int = 25):
    with DB_LOCK:
        return cur.execute(
            "SELECT * FROM ai_chats WHERE user_id=? AND (closed_at IS NULL OR closed_at='') ORDER BY id DESC LIMIT ?",
            (user_id, limit)
        ).fetchall()

def add_msg(ai_chat_id: int, role: str, content: str):
    with DB_LOCK:
        cur.execute(
            "INSERT INTO ai_messages(ai_chat_id, role, content, created_at) VALUES (?,?,?,?)",
            (ai_chat_id, role, str(content), now_str())
        )
        conn.commit()

def get_recent_messages(ai_chat_id: int, limit: int = 16):
    with DB_LOCK:
        rows = cur.execute(
            "SELECT role, content FROM ai_messages WHERE ai_chat_id=? ORDER BY id DESC LIMIT ?",
            (ai_chat_id, limit)
        ).fetchall()
    rows = list(reversed(rows))
    return [{"role": r["role"], "content": r["content"]} for r in rows]


# ================== BO‘LIMLAR / STATES ==================
SECTION = {}  # "home" | "ai" | "media" | "pptx"
WAIT_RENAME_ACTIVE_CHAT = set()

# Media stream
MEDIA_MODE = {}          # user_id -> None | "img" | "vid"
MEDIA_QUERY = {}         # user_id -> query
MEDIA_USED = {}          # user_id -> set(url/link)

# PPTX flow
PPTX_FLOW = {}           # user_id -> dict


# ================== UI TEXTS ==================
BTN_HOME = "🏠 Bosh menyu"

BTN_AI_SECTION = "🤖 AI bo‘limi"
BTN_MEDIA_SECTION = "🧰 Media bo‘limi"
BTN_PPTX_SECTION = "📊 PPTX bo‘limi"

# AI buttons
BTN_NEW_CHAT = "🧠 Yangi chat"
BTN_CHAT_LIST = "🗂 Chatlarim"
BTN_RENAME_CHAT = "✏️ Chat nomini o‘zgartirish"

# Media buttons
BTN_IMG = "🖼 Rasm qidirish"
BTN_VID = "🎬 Video qidirish"
BTN_MEDIA_STOP = "🛑 Media to‘xtatish"
BTN_BACK_AI = "↩️ AI ga qaytish"

# PPTX buttons
BTN_PPTX_START = "📊 PPTX yaratish"
BTN_PPTX_CANCEL = "🛑 PPTX bekor qilish"


# ================== MENUS ==================
def menu_home():
    kb = ReplyKeyboardMarkup(resize_keyboard=True)
    kb.row(KeyboardButton(BTN_AI_SECTION), KeyboardButton(BTN_MEDIA_SECTION))
    kb.row(KeyboardButton(BTN_PPTX_SECTION))
    return kb

def menu_ai():
    kb = ReplyKeyboardMarkup(resize_keyboard=True)
    kb.row(KeyboardButton(BTN_NEW_CHAT), KeyboardButton(BTN_CHAT_LIST))
    kb.row(KeyboardButton(BTN_RENAME_CHAT))
    kb.row(KeyboardButton(BTN_HOME))
    return kb

def menu_media():
    kb = ReplyKeyboardMarkup(resize_keyboard=True)
    kb.row(KeyboardButton(BTN_IMG), KeyboardButton(BTN_VID))
    kb.row(KeyboardButton(BTN_MEDIA_STOP))
    kb.row(KeyboardButton(BTN_BACK_AI), KeyboardButton(BTN_HOME))
    return kb

def menu_pptx():
    kb = ReplyKeyboardMarkup(resize_keyboard=True)
    kb.row(KeyboardButton(BTN_PPTX_START), KeyboardButton(BTN_PPTX_CANCEL))
    kb.row(KeyboardButton(BTN_HOME))
    return kb

def set_section(uid: int, section: str):
    SECTION[uid] = section

def get_section(uid: int) -> str:
    return SECTION.get(uid, "home")

def media_reset(uid: int):
    MEDIA_MODE[uid] = None
    MEDIA_QUERY.pop(uid, None)
    MEDIA_USED[uid] = set()

def reset_states(uid: int):
    WAIT_RENAME_ACTIVE_CHAT.discard(uid)
    PPTX_FLOW.pop(uid, None)
    media_reset(uid)


# ================== MEDIA SEARCH ==================
ALLOWED_IMG_DOMAINS = (
    "upload.wikimedia.org",
    "images.unsplash.com",
    "cdn.pixabay.com",
    "images.pexels.com",
    "i.imgur.com",
)

def _host(url: str) -> str:
    try:
        return url.split("//", 1)[1].split("/", 1)[0].lower()
    except Exception:
        return ""

def find_image_urls_ddg(query: str, max_results: int = 25):
    with DDGS() as ddgs:
        results = ddgs.images(
            query,
            max_results=max_results,
            safesearch="moderate",
            size="Large"
        )
    urls = []
    for r in results:
        u = r.get("image")
        if not u or not isinstance(u, str) or not u.startswith("http"):
            continue
        h = _host(u)
        if any(h.endswith(d) for d in ALLOWED_IMG_DOMAINS):
            urls.append(u)

    # agar filtrdan hech narsa chiqmasa, filtrsiz ham ozgina qaytaramiz
    if not urls:
        for r in results:
            u = r.get("image")
            if u and isinstance(u, str) and u.startswith("http"):
                urls.append(u)

    random.shuffle(urls)
    return urls

def ddg_video_search(query: str, max_results: int = 10):
    q = f"{query} video"
    out = []
    with DDGS() as ddgs:
        for r in ddgs.text(q, max_results=max_results, safesearch="moderate"):
            title = (r.get("title") or "Video")
            href = r.get("href") or r.get("url")
            if href and isinstance(href, str) and href.startswith("http"):
                out.append((title, href))
    return out

def download_image_bytes(url: str):
    headers = {
        "User-Agent": "Mozilla/5.0",
        "Accept": "image/*,*/*;q=0.8",
        "Referer": "https://duckduckgo.com/"
    }

    # avval HEAD bilan tekshiramiz
    try:
        h = session.head(url, headers=headers, timeout=10, allow_redirects=True, proxies=NO_PROXY)
        ctype = (h.headers.get("Content-Type") or "").lower()
        if ctype and ("image" not in ctype):
            return None
    except Exception:
        # HEAD ishlamasa ham urinib ko‘ramiz
        pass

    resp = session.get(url, headers=headers, timeout=20, proxies=NO_PROXY, allow_redirects=True)
    resp.raise_for_status()
    ctype = (resp.headers.get("Content-Type") or "").lower()
    if "image" not in ctype:
        return None
    return resp.content


# ================== VISION HELPERS ==================
def file_to_data_url(path: str) -> str:
    mime, _ = mimetypes.guess_type(path)
    if not mime:
        mime = "image/jpeg"
    with open(path, "rb") as f:
        b64 = base64.b64encode(f.read()).decode("utf-8")
    return f"data:{mime};base64,{b64}"


# ================== PPTX (LOCAL TEMPLATES) ==================
# templates/ ichidagi fayllaringiz nomiga moslab yozing
TEMPLATES = {
    "tech_dark": [
        {"key": "tech_dark", "name": "Tech Dark", "file": "templates/tech_dark.pptx", "preview": None},
    ],
    "business": [
        {"key": "business", "name": "Business", "file": "templates/business.pptx", "preview": None},
    ],
    "education": [
        {"key": "school_smiles", "name": "School Smiles", "file": "templates/school_smiles.pptx", "preview": None},
    ],
    "minimal": [
        {"key": "minimal", "name": "Minimal", "file": "templates/minimal.pptx", "preview": None},
    ],
}

STYLE_LABELS = {
    "tech_dark": "🌙 Tech / Dark",
    "business": "💼 Biznes",
    "education": "🎓 Ta’lim",
    "minimal": "✨ Minimal",
}

def _delete_all_slides(prs: Presentation):
    slide_ids = list(prs.slides._sldIdLst)
    for sldId in slide_ids:
        prs.slides._sldIdLst.remove(sldId)

def _pick_layout(prs: Presentation):
    # template’larda nomlar har xil bo‘lishi mumkin
    prefer = ("Title and Content", "Title & Content", "Content", "Two Content")
    for name in prefer:
        for layout in prs.slide_layouts:
            if (layout.name or "").lower().find(name.lower()) != -1:
                return layout
    return prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

def _set_title(slide, title: str):
    try:
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            slide.shapes.title.text_frame.clear()
            slide.shapes.title.text_frame.paragraphs[0].text = title
            return True
    except Exception:
        pass

    # fallback: birinchi text_frame
    for sh in slide.shapes:
        if sh.has_text_frame:
            sh.text_frame.clear()
            sh.text_frame.paragraphs[0].text = title
            return True
    return False

def _find_best_body_shape(slide):
    best = None
    best_area = -1

    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        try:
            if slide.shapes.title and sh == slide.shapes.title:
                continue
        except Exception:
            pass

        try:
            area = int(sh.width) * int(sh.height)
        except Exception:
            area = 0

        if area > best_area:
            best_area = area
            best = sh

    return best

def _set_bullets(slide, bullets: list[str]):
    body = _find_best_body_shape(slide)
    if not body:
        # fallback: hech narsa topilmasa title’ning o‘ziga yozib yuboramiz
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                tf = slide.shapes.title.text_frame
                tf.text = (slide.shapes.title.text + "\n" + "\n".join(bullets[:6])).strip()
                return True
        except Exception:
            return False

    tf = body.text_frame
    tf.clear()

    if not bullets:
        bullets = ["Ma’lumot"]

    p0 = tf.paragraphs[0]
    p0.text = str(bullets[0])
    p0.level = 0

    for b in bullets[1:8]:
        p = tf.add_paragraph()
        p.text = str(b)
        p.level = 0
    return True

def groq_make_slide_plan(topic: str, slides_count: int) -> list[dict]:
    slides_count = max(1, min(20, int(slides_count)))

    prompt = (
        "Menga PowerPoint prezentatsiya uchun slayd rejasi kerak.\n"
        f"Mavzu: {topic}\n"
        f"Slayd soni: {slides_count}\n\n"
        "FAKAT JSON qaytar:\n"
        "{\n"
        '  "slides": [\n'
        '     {"title": "...", "bullets": ["...","...","..."]}\n'
        "  ]\n"
        "}\n\n"
        "Qoidalar:\n"
        "- slides uzunligi aynan ko‘rsatilgan slayd soniga teng bo‘lsin\n"
        "- matnlar o‘zbek tilida bo‘lsin\n"
        "- har slaydda 3-6 bullet bo‘lsin\n"
        "- juda uzun gap yozmang\n"
    )

    res = client.chat.completions.create(
        model=CHAT_MODEL,
        messages=[
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user", "content": prompt}
        ],
        temperature=0.4
    )

    raw = (res.choices[0].message.content or "").strip()
    m = re.search(r"\{[\s\S]*\}", raw)
    if m:
        raw = m.group(0)

    try:
        data = json.loads(raw)
        slides = data.get("slides") or []
    except Exception:
        slides = []

    out = []
    for s in slides[:slides_count]:
        t = str(s.get("title") or "").strip() or "Slayd"
        b = s.get("bullets") or []
        b = [str(x).strip() for x in b if str(x).strip()]
        if len(b) < 3:
            b = (b + ["Ma’lumot"] * 3)[:3]
        out.append({"title": t, "bullets": b[:8]})

    while len(out) < slides_count:
        out.append({"title": f"Slayd {len(out)+1}", "bullets": ["Ma’lumot", "Misol", "Xulosa"]})

    return out

def build_pptx_from_template(template_path: str, topic: str, slides: list[dict]) -> str:
    if not os.path.exists(template_path):
        raise RuntimeError(f"Template fayl topilmadi: {template_path}")

    # Duplicate warnings ko‘rinmasin
    warnings.filterwarnings("ignore", message=r"Duplicate name: 'ppt/slides/.*'")
    warnings.filterwarnings("ignore", message=r"Duplicate name: 'ppt/slides/_rels/.*'")

    prs = Presentation(template_path)
    _delete_all_slides(prs)

    title_layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else prs.slide_layouts[0]
    body_layout = _pick_layout(prs)

    # Title slide
    s0 = prs.slides.add_slide(title_layout)
    _set_title(s0, topic)
    _set_bullets(s0, [""])

    # Content slides
    for s in slides[1:]:
        sl = prs.slides.add_slide(body_layout)
        _set_title(sl, s["title"])
        _set_bullets(sl, s["bullets"])

    out_name = f"pptx_out/{uuid.uuid4().hex[:10]}_{int(time.time())}.pptx"
    prs.save(out_name)
    return out_name


# ================== CALLBACKS (AI CHAT) ==================
@bot.callback_query_handler(func=lambda c: c.data in ["ai_newchat_ok", "ai_newchat_cancel"])
def ai_newchat_cb(call):
    uid = call.from_user.id
    if get_section(uid) != "ai":
        bot.answer_callback_query(call.id, "AI bo‘limida emassiz")
        return

    if call.data == "ai_newchat_cancel":
        bot.answer_callback_query(call.id, "Bekor qilindi")
        return

    active = get_active_chat_id(uid)
    if active:
        close_chat(uid, active)

    new_name = f"Chat {int(time.time())}"
    create_ai_chat(uid, new_name)

    bot.answer_callback_query(call.id, "Yangi chat ✅")
    send_html(call.message.chat.id, f"✅ Yangi chat ochildi: <b>{safe_html(new_name)}</b>", reply_markup=menu_ai())

@bot.callback_query_handler(func=lambda c: c.data.startswith("ai_switch:"))
def ai_switch_cb(call):
    uid = call.from_user.id
    if get_section(uid) != "ai":
        bot.answer_callback_query(call.id, "AI bo‘limida emassiz")
        return

    chat_id = int(call.data.split(":", 1)[1])
    with DB_LOCK:
        row = cur.execute(
            "SELECT id FROM ai_chats WHERE id=? AND user_id=? AND (closed_at IS NULL OR closed_at='')",
            (chat_id, uid)
        ).fetchone()

    if not row:
        bot.answer_callback_query(call.id, "Bu chat sizniki emas ❌", show_alert=True)
        return

    set_active_chat(uid, chat_id)
    bot.answer_callback_query(call.id, "Aktiv chat ✅")
    send_html(call.message.chat.id, "✅ Aktiv chat o‘zgardi.", reply_markup=menu_ai())

@bot.callback_query_handler(func=lambda c: c.data == "ai_list_close")
def ai_list_close(call):
    bot.answer_callback_query(call.id, "Yopildi")
    try:
        bot.delete_message(call.message.chat.id, call.message.message_id)
    except Exception:
        pass


# ================== CALLBACKS (PPTX) ==================
@bot.callback_query_handler(func=lambda c: c.data.startswith("pptx_style:"))
def pptx_style_cb(call):
    uid = call.from_user.id
    if get_section(uid) != "pptx":
        bot.answer_callback_query(call.id, "PPTX bo‘limida emassiz", show_alert=True)
        return

    flow = PPTX_FLOW.get(uid)
    if not flow:
        bot.answer_callback_query(call.id, "Jarayon topilmadi", show_alert=True)
        return

    style = call.data.split(":", 1)[1]
    if style not in TEMPLATES:
        bot.answer_callback_query(call.id, "Stil topilmadi", show_alert=True)
        return

    flow["style"] = style
    flow["step"] = "template"
    bot.answer_callback_query(call.id, "Stil tanlandi ✅")

    send_html(call.message.chat.id, f"🧩 Shablon tanlang: <b>{safe_html(STYLE_LABELS.get(style, style))}</b>")

    for tpl in TEMPLATES[style]:
        kb = InlineKeyboardMarkup()
        kb.add(InlineKeyboardButton("✅ Shu shablonni tanlash", callback_data=f"pptx_tpl:{tpl['key']}"))
        send_plain(call.message.chat.id, f"📌 {tpl['name']}", reply_markup=kb)

@bot.callback_query_handler(func=lambda c: c.data.startswith("pptx_tpl:"))
def pptx_tpl_cb(call):
    uid = call.from_user.id
    if get_section(uid) != "pptx":
        bot.answer_callback_query(call.id, "PPTX bo‘limida emassiz", show_alert=True)
        return

    flow = PPTX_FLOW.get(uid)
    if not flow:
        bot.answer_callback_query(call.id, "Jarayon topilmadi", show_alert=True)
        return

    tpl_key = call.data.split(":", 1)[1]
    style = flow.get("style")

    tpl = None
    for t in TEMPLATES.get(style, []):
        if t["key"] == tpl_key:
            tpl = t
            break
    if not tpl:
        bot.answer_callback_query(call.id, "Shablon topilmadi", show_alert=True)
        return

    topic = flow["topic"]
    slides_count = int(flow["slides"])

    bot.answer_callback_query(call.id, "Yaratilmoqda...")

    wait = send_plain(call.message.chat.id, "⏳ PPTX tayyorlanyapti... (AI + local template)")

    try:
        plan = groq_make_slide_plan(topic, slides_count)
        out_path = build_pptx_from_template(tpl["file"], topic, plan)

        with open(out_path, "rb") as f:
            bot.send_document(call.message.chat.id, f, caption="✅ PPTX tayyor!", parse_mode=None)

        try:
            os.remove(out_path)
        except Exception:
            pass

        PPTX_FLOW.pop(uid, None)
        send_plain(call.message.chat.id, "✅ PPTX tugadi. PPTX bo‘limida qolavering yoki 🏠 Bosh menyuga qayting.")

    except Exception as e:
        send_plain(call.message.chat.id, f"❌ PPTX xato: {e}")
    finally:
        try:
            bot.delete_message(call.message.chat.id, wait.message_id)
        except Exception:
            pass


# ================== /start ==================
@bot.message_handler(commands=["start"])
def start(message):
    uid = message.from_user.id
    chat_id = message.chat.id

    get_or_create_user(uid, chat_id)
    if not get_active_chat_id(uid):
        create_ai_chat(uid, "Chat 1")

    reset_states(uid)
    set_section(uid, "home")

    send_html(chat_id, "🏠 <b>Bosh menyu</b>\n👇 Bo‘lim tanlang:", reply_markup=menu_home())


# ================== HOME NAV ==================
@bot.message_handler(func=lambda m: m.text == BTN_HOME)
def back_home(message):
    uid = message.from_user.id
    reset_states(uid)
    set_section(uid, "home")
    send_html(message.chat.id, "🏠 <b>Bosh menyu</b>\n👇 Bo‘lim tanlang:", reply_markup=menu_home())

@bot.message_handler(func=lambda m: m.text == BTN_AI_SECTION)
def go_ai(message):
    uid = message.from_user.id
    PPTX_FLOW.pop(uid, None)
    media_reset(uid)
    set_section(uid, "ai")
    send_html(message.chat.id, "🤖 <b>AI bo‘limi</b>\nMatn yozsangiz Groq javob beradi.", reply_markup=menu_ai())

@bot.message_handler(func=lambda m: m.text == BTN_MEDIA_SECTION)
def go_media(message):
    uid = message.from_user.id
    PPTX_FLOW.pop(uid, None)
    media_reset(uid)
    set_section(uid, "media")
    send_html(message.chat.id, "🧰 <b>Media bo‘limi</b>\nBu yerda AI javob bermaydi.", reply_markup=menu_media())

@bot.message_handler(func=lambda m: m.text == BTN_PPTX_SECTION)
def go_pptx(message):
    uid = message.from_user.id
    media_reset(uid)
    PPTX_FLOW.pop(uid, None)
    set_section(uid, "pptx")
    send_html(message.chat.id, "📊 <b>PPTX bo‘limi</b>\nBu yerda AI chat javob bermaydi.", reply_markup=menu_pptx())


# ================== AI BUTTONS ==================
@bot.message_handler(func=lambda m: m.text == BTN_NEW_CHAT)
def ai_new_chat_confirm(message):
    uid = message.from_user.id
    if get_section(uid) != "ai":
        return

    kb = InlineKeyboardMarkup()
    kb.add(
        InlineKeyboardButton("✅ Tasdiqlash", callback_data="ai_newchat_ok"),
        InlineKeyboardButton("❌ Bekor qilish", callback_data="ai_newchat_cancel")
    )
    send_html(
        message.chat.id,
        "⚠️ <b>Ishonchingiz komilmi?</b>\n"
        "Bu chatni yopganingizdan keyin bu chatdagi ma’lumotlar qaytib ko‘rsatilmaydi.",
        reply_markup=kb
    )

@bot.message_handler(func=lambda m: m.text == BTN_CHAT_LIST)
def ai_chat_list(message):
    uid = message.from_user.id
    if get_section(uid) != "ai":
        return

    active = get_active_chat_id(uid)
    chats = list_user_chats(uid, limit=25)
    if not chats:
        reply_plain(message, "📭 Chat yo‘q.")
        return

    kb = InlineKeyboardMarkup()
    for c in chats:
        mark = "✅ " if c["id"] == active else ""
        kb.add(InlineKeyboardButton(f"{mark}{c['name']}  (#{c['id']})", callback_data=f"ai_switch:{c['id']}"))
    kb.add(InlineKeyboardButton("❌ Yopish", callback_data="ai_list_close"))

    send_html(message.chat.id, "🗂 <b>Chatlarim</b> (faqat sizniki):", reply_markup=kb)

@bot.message_handler(func=lambda m: m.text == BTN_RENAME_CHAT)
def ai_rename_start(message):
    uid = message.from_user.id
    if get_section(uid) != "ai":
        return
    WAIT_RENAME_ACTIVE_CHAT.add(uid)
    reply_plain(message, "✏️ Yangi chat nomini yozing (40 belgigacha):")

@bot.message_handler(func=lambda m: (m.from_user.id in WAIT_RENAME_ACTIVE_CHAT) and m.text and not m.text.startswith("/"))
def ai_rename_set(message):
    uid = message.from_user.id
    WAIT_RENAME_ACTIVE_CHAT.discard(uid)

    active = get_active_chat_id(uid)
    if not active:
        reply_plain(message, "❌ Aktiv chat topilmadi.")
        return

    nm = (message.text or "").strip()[:40]
    if not nm:
        reply_plain(message, "❌ Nom bo‘sh bo‘lmasin.")
        return

    rename_chat(uid, active, nm)
    reply_html(message, f"✅ Chat nomi o‘zgardi: <b>{safe_html(nm)}</b>")


# ================== MEDIA BUTTONS ==================
@bot.message_handler(func=lambda m: m.text == BTN_BACK_AI)
def media_back_ai(message):
    uid = message.from_user.id
    media_reset(uid)
    set_section(uid, "ai")
    send_html(message.chat.id, "↩️ AI bo‘limiga qaytdingiz.", reply_markup=menu_ai())

@bot.message_handler(func=lambda m: m.text == BTN_IMG)
def media_img_start(message):
    uid = message.from_user.id
    if get_section(uid) != "media":
        return
    MEDIA_MODE[uid] = "img"
    MEDIA_QUERY[uid] = None
    MEDIA_USED[uid] = set()
    reply_plain(message, "🖼 Rasm qidirish tanlandi.\nEndi so‘rov yozing (har xabar -> yangi rasm).")

@bot.message_handler(func=lambda m: m.text == BTN_VID)
def media_vid_start(message):
    uid = message.from_user.id
    if get_section(uid) != "media":
        return
    MEDIA_MODE[uid] = "vid"
    MEDIA_QUERY[uid] = None
    MEDIA_USED[uid] = set()
    reply_plain(message, "🎬 Video qidirish tanlandi.\nEndi so‘rov yozing (har xabar -> yangi linklar).")

@bot.message_handler(func=lambda m: m.text == BTN_MEDIA_STOP)
def media_stop(message):
    uid = message.from_user.id
    if get_section(uid) != "media":
        return
    media_reset(uid)
    reply_plain(message, "🛑 Media qidirish to‘xtatildi. Yana 🖼/🎬 tanlashingiz mumkin.")


# ================== PPTX FLOW ==================
@bot.message_handler(func=lambda m: m.text == BTN_PPTX_START)
def pptx_start(message):
    uid = message.from_user.id
    if get_section(uid) != "pptx":
        return
    PPTX_FLOW[uid] = {"step": "topic"}
    reply_html(message, "📊 PPTX yaratamiz.\n\n1) Mavzuni yozing (masalan: <b>Ekologik muammolar</b>)")

@bot.message_handler(func=lambda m: m.text == BTN_PPTX_CANCEL)
def pptx_cancel(message):
    uid = message.from_user.id
    if get_section(uid) != "pptx":
        return
    PPTX_FLOW.pop(uid, None)
    reply_plain(message, "🛑 PPTX bekor qilindi.")


# ================== VOICE (WHISPER) ==================
@bot.message_handler(content_types=["voice"])
def handle_voice(message):
    uid = message.from_user.id
    get_or_create_user(uid, message.chat.id)

    try:
        f = bot.get_file(message.voice.file_id)
        audio_bytes = bot.download_file(f.file_path)

        ogg_path = f"media/{message.voice.file_id}.ogg"
        with open(ogg_path, "wb") as out:
            out.write(audio_bytes)

        with open(ogg_path, "rb") as audio_file:
            tr = client.audio.transcriptions.create(
                file=audio_file,
                model=WHISPER_MODEL,
                language="uz",
                response_format="json"
            )

        text = tr.text if hasattr(tr, "text") else tr["text"]
        send_html(message.chat.id, f"🎤 <b>Ovozdan matn:</b>\n{safe_html(text)}")

        # AI bo‘limida bo‘lsa transkriptni AI’ga ham yuboramiz
        if get_section(uid) == "ai":
            active = get_active_chat_id(uid) or create_ai_chat(uid, "Chat 1")
            history = get_recent_messages(active, limit=16)
            msgs = [{"role": "system", "content": SYSTEM_PROMPT}] + history + [{"role": "user", "content": text}]
            add_msg(active, "user", text)

            res = client.chat.completions.create(
                model=CHAT_MODEL,
                messages=msgs,
                temperature=0.3
            )
            ans = res.choices[0].message.content
            add_msg(active, "assistant", ans)
            send_plain(message.chat.id, ans)

    except Exception as e:
        safe_error(message.chat.id, "❌ Ovoz xatosi", e)


# ================== PHOTO (VISION) ==================
@bot.message_handler(content_types=["photo"])
def handle_photo(message):
    uid = message.from_user.id
    get_or_create_user(uid, message.chat.id)

    try:
        file_id = message.photo[-1].file_id
        f = bot.get_file(file_id)
        img_bytes = bot.download_file(f.file_path)

        path = f"media/{file_id}.jpg"
        with open(path, "wb") as out:
            out.write(img_bytes)

        data_url = file_to_data_url(path)

        res = client.chat.completions.create(
            model=VISION_MODEL,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Bu rasmni o‘zbek tilida qisqa va aniq tahlil qilib ber."},
                        {"type": "image_url", "image_url": {"url": data_url}}
                    ]
                }
            ],
            temperature=0.2
        )

        send_html(message.chat.id, "🖼 <b>Rasm tahlili:</b>\n" + safe_html(res.choices[0].message.content))

    except Exception as e:
        safe_error(message.chat.id, "❌ Rasm xatosi", e)


# ================== MAIN TEXT ROUTER ==================
@bot.message_handler(func=lambda m: m.text and not m.text.startswith("/"))
def router(message):
    uid = message.from_user.id
    chat_id = message.chat.id
    get_or_create_user(uid, chat_id)

    section = get_section(uid)
    txt = (message.text or "").strip()

    # ---------------- PPTX steps ----------------
    if section == "pptx":
        flow = PPTX_FLOW.get(uid)
        if not flow:
            return

        step = flow.get("step")

        if step == "topic":
            if len(txt) < 3:
                reply_plain(message, "Mavzu juda qisqa. Qayta yozing.")
                return
            flow["topic"] = txt
            flow["step"] = "slides"
            reply_plain(message, "2) Nechta slayd bo‘lsin? (1..20)")
            return

        if step == "slides":
            if not txt.isdigit():
                reply_plain(message, "Faqat raqam yuboring (1..20).")
                return
            n = int(txt)
            if n < 1 or n > 20:
                reply_plain(message, "1..20 oralig‘ida bo‘lsin.")
                return
            flow["slides"] = n
            flow["step"] = "style"

            kb = InlineKeyboardMarkup()
            kb.add(
                InlineKeyboardButton(STYLE_LABELS["tech_dark"], callback_data="pptx_style:tech_dark"),
                InlineKeyboardButton(STYLE_LABELS["business"], callback_data="pptx_style:business"),
            )
            kb.add(
                InlineKeyboardButton(STYLE_LABELS["education"], callback_data="pptx_style:education"),
                InlineKeyboardButton(STYLE_LABELS["minimal"], callback_data="pptx_style:minimal"),
            )
            send_html(chat_id, "3) Stil tanlang:", reply_markup=kb)
            return

        return

    # ---------------- MEDIA stream ----------------
    if section == "media":
        mode = MEDIA_MODE.get(uid)
        if mode not in ("img", "vid"):
            return

        query = txt
        if not query:
            return

        if MEDIA_QUERY.get(uid) != query:
            MEDIA_QUERY[uid] = query
            MEDIA_USED[uid] = set()

        if mode == "img":
            try:
                urls = find_image_urls_ddg(query, max_results=18)
                used = MEDIA_USED.get(uid, set())
                last_err = None

                for url in urls:
                    if url in used:
                        continue
                    try:
                        img_bytes = download_image_bytes(url)
                        used.add(url)
                        if not img_bytes:
                            continue
                        MEDIA_USED[uid] = used
                        bot.send_photo(chat_id, img_bytes, caption=f"🖼 {query}", parse_mode=None)
                        return
                    except Exception as e:
                        last_err = e
                        used.add(url)
                        continue

                MEDIA_USED[uid] = used
# ✅ rasmni yubora olmasak ham linklarni tashlab beramiz
                links = []
                for u in urls[:6]:
                    if u not in used:
                        links.append(u)
                    if links:
                        send_plain(chat_id,
                        "⚠️ Rasmni server yuklay olmadi (internet cheklovi).\n"
                        "Lekin linklar:\n\n" + "\n".join(links)
                        )
                    else:
                        send_plain(chat_id, f"❌ Rasm yuborilmadi.\nXato: {last_err}")


            except Exception as e:
                send_plain(chat_id, f"❌ Rasm qidirishda xato: {e}")
                return

        if mode == "vid":
            try:
                res = ddg_video_search(query, max_results=10)
                if not res:
                    send_plain(chat_id, "❌ Video topilmadi.")
                    return

                used = MEDIA_USED.get(uid, set())
                lines = [f"🎬 Topilgan linklar: {query}\n"]
                c = 0
                for title, link in res:
                    if link in used:
                        continue
                    used.add(link)
                    c += 1
                    lines.append(f"{c}) {title}\n{link}\n")
                    if c >= 6:
                        break

                MEDIA_USED[uid] = used
                if c == 0:
                    send_plain(chat_id, "⚠️ Yangi link topilmadi. Boshqa so‘z bilan urinib ko‘ring.")
                    return

                send_plain(chat_id, "\n".join(lines))
                return

            except Exception as e:
                send_plain(chat_id, f"❌ Video qidirishda xato: {e}")
                return

    # ---------------- AI section ----------------
    if section != "ai":
        return

    if uid in WAIT_RENAME_ACTIVE_CHAT:
        return

    # tugmalarni yeb yubormasin
    if txt in (
        BTN_NEW_CHAT, BTN_CHAT_LIST, BTN_RENAME_CHAT,
        BTN_AI_SECTION, BTN_MEDIA_SECTION, BTN_PPTX_SECTION,
        BTN_IMG, BTN_VID, BTN_MEDIA_STOP, BTN_BACK_AI,
        BTN_PPTX_START, BTN_PPTX_CANCEL, BTN_HOME
    ):
        return

    active = get_active_chat_id(uid) or create_ai_chat(uid, "Chat 1")

    try:
        history = get_recent_messages(active, limit=16)
        msgs = [{"role": "system", "content": SYSTEM_PROMPT}] + history + [{"role": "user", "content": txt}]
        add_msg(active, "user", txt)

        res = client.chat.completions.create(
            model=CHAT_MODEL,
            messages=msgs,
            temperature=0.3
        )
        ans = res.choices[0].message.content
        add_msg(active, "assistant", ans)
        send_plain(chat_id, ans)

    except Exception as e:
        send_plain(chat_id, f"❌ AI xato: {e}")


# ================== BO‘LIM KIRISH TUGMALARI ==================
@bot.message_handler(func=lambda m: m.text == BTN_BACK_AI)
def back_ai(message):
    uid = message.from_user.id
    media_reset(uid)
    set_section(uid, "ai")
    send_html(message.chat.id, "↩️ AI bo‘limiga qaytdingiz.", reply_markup=menu_ai())


print("✅ Bot ishga tushdi")
bot.infinity_polling()
